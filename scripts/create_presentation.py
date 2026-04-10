#!/usr/bin/env python3
"""
Script to create AI Productivity Presentation PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
GREEN = RGBColor(0, 176, 80)
RED = RGBColor(255, 0, 0)
ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(89, 89, 89)


def add_title_slide(prs):
    """Add title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI-Powered SIEM Extension Development"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Productivity Enhancement Report"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "March 2026"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER


def add_before_ai_slide(prs):
    """Slide 1: Before AI Workflow"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Before AI: Traditional Workflow"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Workflow steps - horizontal flow
    steps = [
        ("1", "Receive\nRequest", LIGHT_BLUE),
        ("2", "Manual Log\nAnalysis", LIGHT_BLUE),
        ("3", "Schema\nVerification", LIGHT_BLUE),
        ("4", "Compare with\nVendor Docs", LIGHT_BLUE),
        ("5", "Manual Use Case\nDevelopment", LIGHT_BLUE),
        ("6", "Testing &\nDeploy", LIGHT_BLUE),
    ]
    
    start_x = 0.3
    for i, (num, text, color) in enumerate(steps):
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_x + i * 2.1), Inches(1.5), Inches(1.9), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = DARK_BLUE
        
        # Text in box
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{num}. {text}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(8)
        
        # Arrow (except last)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(start_x + i * 2.1 + 1.95), Inches(1.85), Inches(0.2), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()
    
    # Time breakdown table
    table_data = [
        ["Step", "Activity", "Time", "Pain Points"],
        ["1", "Log Analysis", "2-3 hours", "Manual parsing, pattern identification"],
        ["2", "Vendor Doc Review", "1-2 hours", "Cross-referencing, schema mapping"],
        ["3", "Schema Verification", "1-2 hours", "Field-by-field comparison"],
        ["4", "Use Case Development", "3-4 hours", "Manual rule creation, testing iterations"],
        ["5", "Documentation", "1 hour", "Manual report writing"],
    ]
    
    rows, cols = len(table_data), len(table_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(2.8), Inches(9), Inches(2.5)).table
    
    col_widths = [Inches(0.6), Inches(2.2), Inches(1.2), Inches(5)]
    for i, width in enumerate(col_widths):
        table.columns[i].width = width
    
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.alignment = PP_ALIGN.CENTER
            if row_idx == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            else:
                para.font.color.rgb = GRAY
    
    # Total time box
    total_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(2.8), Inches(3.5), Inches(1))
    total_box.fill.solid()
    total_box.fill.fore_color.rgb = RED
    total_box.line.fill.background()
    tf = total_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Total Time:\n8-12 hours"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Bottlenecks
    bottleneck_title = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(3), Inches(0.5))
    tf = bottleneck_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Bottlenecks:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    bottlenecks = [
        "❌ Repetitive manual analysis for each data source",
        "❌ High cognitive load for schema comparison",
        "❌ Error-prone manual rule development",
        "❌ Limited scalability with increasing requests"
    ]
    
    for i, text in enumerate(bottlenecks):
        box = slide.shapes.add_textbox(Inches(0.3 + (i % 2) * 6.5), Inches(5.9 + (i // 2) * 0.4), Inches(6), Inches(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY


def add_after_ai_slide(prs):
    """Slide 2: After AI Workflow"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = GREEN
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "After AI: Streamlined Workflow"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # New workflow - visual flow
    steps = [
        ("Vendor Docs +\nSample Logs", LIGHT_BLUE),
        ("AI Analysis", GREEN),
        ("Alert Suggestions +\nImportant Fields", GREEN),
        ("Manual\nVerification", ORANGE),
        ("Final\nReport", LIGHT_BLUE),
    ]
    
    start_x = 0.5
    for i, (text, color) in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_x + i * 2.5), Inches(1.5), Inches(2.2), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = DARK_BLUE
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(12)
        
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(start_x + i * 2.5 + 2.25), Inches(1.9), Inches(0.3), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()
    
    # Optimized time table
    table_data = [
        ["Step", "Activity", "Time", "AI Contribution"],
        ["1", "Input Collection", "15-30 min", "Vendor docs + sample logs upload"],
        ["2", "AI Analysis", "5-10 min", "Automated schema & pattern analysis"],
        ["3", "Review Suggestions", "30-45 min", "AI-generated alerts & field mapping"],
        ["4", "Manual Verification", "1-2 hours", "Human validation & fine-tuning"],
        ["5", "Report Generation", "15-30 min", "AI-assisted documentation"],
    ]
    
    rows, cols = len(table_data), len(table_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(2.9), Inches(9), Inches(2.5)).table
    
    col_widths = [Inches(0.6), Inches(2), Inches(1.3), Inches(5.1)]
    for i, width in enumerate(col_widths):
        table.columns[i].width = width
    
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.alignment = PP_ALIGN.CENTER
            if row_idx == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN
            else:
                para.font.color.rgb = GRAY
    
    # Total time box
    total_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(2.9), Inches(3.5), Inches(1))
    total_box.fill.solid()
    total_box.fill.fore_color.rgb = GREEN
    total_box.line.fill.background()
    tf = total_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Total Time:\n4-5 hours"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Benefits
    benefits_title = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(3), Inches(0.5))
    tf = benefits_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Benefits:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    benefits = [
        "✅ Automated log parsing and schema comparison",
        "✅ AI-suggested use cases based on best practices",
        "✅ Consistent field identification across data sources",
        "✅ Reduced human error in rule development"
    ]
    
    for i, text in enumerate(benefits):
        box = slide.shapes.add_textbox(Inches(0.3 + (i % 2) * 6.5), Inches(5.9 + (i // 2) * 0.4), Inches(6), Inches(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY


def add_comparison_slide(prs):
    """Slide 3: Productivity Comparison"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Productivity Comparison: Before vs After AI"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Time per Extension comparison
    label1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(3), Inches(0.4))
    tf = label1.text_frame
    p = tf.paragraphs[0]
    p.text = "Time per Extension"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Before bar
    before_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.95), Inches(1.5), Inches(0.4))
    tf = before_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Before AI"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    
    before_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.9), Inches(8), Inches(0.4))
    before_bar.fill.solid()
    before_bar.fill.fore_color.rgb = RED
    before_bar.line.fill.background()
    
    before_time = slide.shapes.add_textbox(Inches(10.1), Inches(1.95), Inches(2), Inches(0.4))
    tf = before_time.text_frame
    p = tf.paragraphs[0]
    p.text = "8-12 hours"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RED
    
    # After bar
    after_label = slide.shapes.add_textbox(Inches(0.5), Inches(2.45), Inches(1.5), Inches(0.4))
    tf = after_label.text_frame
    p = tf.paragraphs[0]
    p.text = "After AI"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    
    after_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.4), Inches(4), Inches(0.4))
    after_bar.fill.solid()
    after_bar.fill.fore_color.rgb = GREEN
    after_bar.line.fill.background()
    
    after_time = slide.shapes.add_textbox(Inches(6.1), Inches(2.45), Inches(2), Inches(0.4))
    tf = after_time.text_frame
    p = tf.paragraphs[0]
    p.text = "4-5 hours"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    # Weekly Capacity
    label2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(4), Inches(0.4))
    tf = label2.text_frame
    p = tf.paragraphs[0]
    p.text = "Weekly Capacity (40-hour week)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Before capacity
    before_cap_label = slide.shapes.add_textbox(Inches(0.5), Inches(3.55), Inches(1.5), Inches(0.4))
    tf = before_cap_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Before AI"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    
    before_cap_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.5), Inches(3), Inches(0.4))
    before_cap_bar.fill.solid()
    before_cap_bar.fill.fore_color.rgb = RED
    before_cap_bar.line.fill.background()
    
    before_cap = slide.shapes.add_textbox(Inches(5.1), Inches(3.55), Inches(3), Inches(0.4))
    tf = before_cap.text_frame
    p = tf.paragraphs[0]
    p.text = "3-5 extensions/week"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RED
    
    # After capacity
    after_cap_label = slide.shapes.add_textbox(Inches(0.5), Inches(4.05), Inches(1.5), Inches(0.4))
    tf = after_cap_label.text_frame
    p = tf.paragraphs[0]
    p.text = "After AI"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    
    after_cap_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(4), Inches(6), Inches(0.4))
    after_cap_bar.fill.solid()
    after_cap_bar.fill.fore_color.rgb = GREEN
    after_cap_bar.line.fill.background()
    
    after_cap = slide.shapes.add_textbox(Inches(8.1), Inches(4.05), Inches(3), Inches(0.4))
    tf = after_cap.text_frame
    p = tf.paragraphs[0]
    p.text = "8-10 extensions/week"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    # Big productivity gain box
    gain_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(4.7), Inches(5.3), Inches(1.3))
    gain_box.fill.solid()
    gain_box.fill.fore_color.rgb = GREEN
    gain_box.line.fill.background()
    tf = gain_box.text_frame
    p = tf.paragraphs[0]
    p.text = "40-50% PRODUCTIVITY GAIN"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Time Saved: ~4-7 hours per task"
    p2.font.size = Pt(16)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    # Key metrics table
    table_data = [
        ["Metric", "Before AI", "After AI", "Improvement"],
        ["Time per Extension", "8-12 hours", "4-5 hours", "50-58% faster"],
        ["Extensions per Week", "3-5", "8-10", "2x throughput"],
        ["Error Rate", "Higher", "Lower", "Improved accuracy"],
        ["Consistency", "Variable", "Standardized", "Better quality"],
    ]
    
    rows, cols = len(table_data), len(table_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.3)).table
    
    col_widths = [Inches(3), Inches(3), Inches(3), Inches(3.3)]
    for i, width in enumerate(col_widths):
        table.columns[i].width = width
    
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.alignment = PP_ALIGN.CENTER
            if row_idx == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif col_idx == 3:
                para.font.bold = True
                para.font.color.rgb = GREEN
            else:
                para.font.color.rgb = GRAY


def add_summary_slide(prs):
    """Slide 4: Summary & Impact"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Summary & Business Impact"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Before box
    before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.5), Inches(2.3))
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = RGBColor(255, 230, 230)
    before_box.line.color.rgb = RED
    
    before_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5), Inches(0.5))
    tf = before_title.text_frame
    p = tf.paragraphs[0]
    p.text = "BEFORE AI"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED
    
    before_items = ["• Manual Analysis", "• Time-Intensive", "• Error-Prone", "• Limited Scale", "• 8-12 hrs/task"]
    before_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(5), Inches(1.6))
    tf = before_text.text_frame
    for i, item in enumerate(before_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
    
    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.2), Inches(2.3), Inches(0.9), Inches(0.6))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GREEN
    arrow.line.fill.background()
    
    # After box
    after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.5), Inches(5.5), Inches(2.3))
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
    after_box.line.color.rgb = GREEN
    
    after_title = slide.shapes.add_textbox(Inches(7.5), Inches(1.6), Inches(5), Inches(0.5))
    tf = after_title.text_frame
    p = tf.paragraphs[0]
    p.text = "AFTER AI"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    after_items = ["• AI-Assisted", "• Streamlined", "• Consistent", "• Scalable", "• 4-5 hrs/task"]
    after_text = slide.shapes.add_textbox(Inches(7.5), Inches(2.1), Inches(5), Inches(1.6))
    tf = after_text.text_frame
    for i, item in enumerate(after_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
    
    # Business Impact section
    impact_title = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(6), Inches(0.5))
    tf = impact_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Quantitative Gains"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    quant_items = [
        "• 40-50% Productivity Improvement in extension development",
        "• 2x Throughput - Handle more requests in the same timeframe",
        "• Faster Time-to-Value - Quicker deployment of security coverage"
    ]
    quant_text = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(6), Inches(1.2))
    tf = quant_text.text_frame
    for i, item in enumerate(quant_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
    
    # Qualitative benefits
    qual_title = slide.shapes.add_textbox(Inches(7), Inches(4), Inches(6), Inches(0.5))
    tf = qual_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Qualitative Benefits"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    qual_items = [
        "• Standardization - Consistent approach across data sources",
        "• Knowledge Capture - AI learns from best practices",
        "• Reduced Burnout - Less repetitive manual work",
        "• Focus Shift - More time for strategic work"
    ]
    qual_text = slide.shapes.add_textbox(Inches(7), Inches(4.4), Inches(6), Inches(1.4))
    tf = qual_text.text_frame
    for i, item in enumerate(qual_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
    
    # Human in the loop box
    hitl_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.9))
    hitl_box.fill.solid()
    hitl_box.fill.fore_color.rgb = RGBColor(240, 240, 255)
    hitl_box.line.color.rgb = LIGHT_BLUE
    
    hitl_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(12), Inches(0.7))
    tf = hitl_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Human-in-the-Loop Value:  AI handles schema analysis, pattern detection, initial suggestions  |  "
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    run = p.add_run()
    run.text = "Human adds domain expertise, verification, final approval"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    
    # Conclusion box
    conclusion_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.6))
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = GREEN
    conclusion_box.line.fill.background()
    
    conclusion_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(12), Inches(0.5))
    tf = conclusion_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Workflow: Vendor Docs + Sample Logs → AI Analysis → Alert Suggestions → Manual Verification → Report"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


# Build the presentation
add_title_slide(prs)
add_before_ai_slide(prs)
add_after_ai_slide(prs)
add_comparison_slide(prs)
add_summary_slide(prs)

# Save
output_path = "/Users/vaibhav.tiwari/siem-optimizer/web3-xdr/docs/AI_Productivity_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
